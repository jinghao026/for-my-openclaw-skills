#!/usr/bin/env python3
"""
Document Parser - 通用文档解析工具
支持格式: DOCX, PPTX, XLSX, PDF
"""

import sys
import os

def read_docx(file_path):
    """读取 Word 文档"""
    try:
        from docx import Document
        doc = Document(file_path)
        text = []
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
        return '\n'.join(text)
    except Exception as e:
        return f"Error reading DOCX: {e}"

def read_pptx(file_path):
    """读取 PowerPoint 文档"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = []
        for i, slide in enumerate(prs.slides, 1):
            text.append(f"\n--- Slide {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        return '\n'.join(text)
    except Exception as e:
        return f"Error reading PPTX: {e}"

def read_xlsx(file_path):
    """读取 Excel 文档"""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, data_only=True)
        text = []
        for sheet_name in wb.sheetnames:
            text.append(f"\n=== Sheet: {sheet_name} ===")
            sheet = wb[sheet_name]
            for row in sheet.iter_rows(values_only=True):
                row_text = ' | '.join([str(cell) if cell is not None else '' for cell in row])
                if row_text.strip():
                    text.append(row_text)
        return '\n'.join(text)
    except Exception as e:
        return f"Error reading XLSX: {e}"

def read_pdf(file_path):
    """读取 PDF 文档"""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        text = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            text.append(f"\n--- Page {page_num + 1} ---")
            text.append(page.get_text())
        doc.close()
        return '\n'.join(text)
    except Exception as e:
        return f"Error reading PDF: {e}"

def parse_document(file_path):
    """根据文件扩展名选择解析器"""
    if not os.path.exists(file_path):
        return f"File not found: {file_path}"
    
    ext = os.path.splitext(file_path)[1].lower()
    
    parsers = {
        '.docx': read_docx,
        '.pptx': read_pptx,
        '.xlsx': read_xlsx,
        '.pdf': read_pdf,
    }
    
    if ext in parsers:
        return parsers[ext](file_path)
    else:
        return f"Unsupported file format: {ext}\nSupported formats: {', '.join(parsers.keys())}"

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python document_parser.py <file_path>")
        print("Supported formats: .docx, .pptx, .xlsx, .pdf")
        sys.exit(1)
    
    file_path = sys.argv[1]
    result = parse_document(file_path)
    print(result)
